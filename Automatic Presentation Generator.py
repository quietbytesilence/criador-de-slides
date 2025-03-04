import os, platform, subprocess
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ("left", "center", "right") --> PP_ALIGN.LEFT, PP_ALIGN.CENTER, e PP_ALIGN.RIGHT.

# Variáveis de personalização
COR_FUNDO = RGBColor(0, 0, 0)  # Preto
COR_TEXTO = RGBColor(255, 255, 255)  # Branco

NOME_PROFESSOR = "Prof. Esp. Gustavo Carneiro"

FONTE_TITULO_PRINCIPAL = "Times New Roman"
TAMANHO_TITULO_PRINCIPAL = Pt(44)
ALINHAMENTO_TITULO_PRINCIPAL = PP_ALIGN.CENTER

FONTE_SUBTITULO = "Times New Roman"
TAMANHO_SUBTITULO = Pt(28)
ALINHAMENTO_SUBTITULO = PP_ALIGN.CENTER

FONTE_TITULO_SLIDE = "Times New Roman"
TAMANHO_TITULO_SLIDE = Pt(28)
ALINHAMENTO_TITULO_SLIDE = PP_ALIGN.LEFT

FONTE_CORPO_TEXTO = "Calibri"
TAMANHO_CORPO_TEXTO = Pt(20)
ALINHAMENTO_CORPO_TEXTO = PP_ALIGN.LEFT

def criar_slides_a_partir_de_txt(nome_arquivo_txt):
    # Ler o conteúdo do arquivo de texto
    with open(nome_arquivo_txt, 'r', encoding='utf-8') as file:
        conteudo = file.read()

    # Dividir o conteúdo em parágrafos (ou blocos de texto)
    blocos = conteudo.split('\n\n')  # Assumindo que os parágrafos são separados por duas quebras de linha

    # Criar uma apresentação
    prs = Presentation()

    # Definir o tema escuro (fundo preto e texto branco)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    background = prs.slide_master.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COR_FUNDO  # Fundo preto

    # Adicionar um slide de título
    slide_titulo = prs.slides.add_slide(prs.slide_layouts[0])
    titulo = slide_titulo.shapes.title
    subtitulo = slide_titulo.placeholders[1]
    titulo.text = os.path.splitext(nome_arquivo_txt)[0]  # Título é o nome do arquivo sem extensão
    subtitulo.text = NOME_PROFESSOR

    # Aplicar estilo ao título e subtítulo
    titulo.text_frame.paragraphs[0].font.color.rgb = COR_TEXTO
    titulo.text_frame.paragraphs[0].font.name = FONTE_TITULO_PRINCIPAL
    titulo.text_frame.paragraphs[0].font.size = TAMANHO_TITULO_PRINCIPAL
    titulo.text_frame.paragraphs[0].alignment = ALINHAMENTO_TITULO_PRINCIPAL

    subtitulo.text_frame.paragraphs[0].font.color.rgb = COR_TEXTO
    subtitulo.text_frame.paragraphs[0].font.name = FONTE_SUBTITULO
    subtitulo.text_frame.paragraphs[0].font.size = TAMANHO_SUBTITULO
    subtitulo.text_frame.paragraphs[0].alignment = ALINHAMENTO_SUBTITULO

    # Adicionar slides com os blocos de conteúdo
    for bloco in blocos:
        if bloco.strip():  # Ignorar blocos vazios
            # Dividir o bloco em título e corpo
            linhas = bloco.strip().split('\n')
            titulo_slide_texto = linhas[0]  # A primeira linha é o título
            corpo_slide_texto = '\n'.join(linhas[1:])  # O restante é o corpo

            # Criar o slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            titulo_slide = slide.shapes.title
            corpo = slide.placeholders[1].text_frame

            # Definir o título do slide
            titulo_slide.text = titulo_slide_texto
            titulo_slide.text_frame.paragraphs[0].font.name = FONTE_TITULO_SLIDE
            titulo_slide.text_frame.paragraphs[0].font.size = TAMANHO_TITULO_SLIDE
            titulo_slide.text_frame.paragraphs[0].font.color.rgb = COR_TEXTO
            titulo_slide.text_frame.paragraphs[0].alignment = ALINHAMENTO_TITULO_SLIDE

            # Adicionar o conteúdo ao corpo do slide
            corpo.text = corpo_slide_texto

            # Aplicar estilo ao corpo do slide
            for paragraph in corpo.paragraphs:
                paragraph.font.name = FONTE_CORPO_TEXTO
                paragraph.font.color.rgb = COR_TEXTO
                paragraph.font.size = TAMANHO_CORPO_TEXTO
                paragraph.alignment = ALINHAMENTO_CORPO_TEXTO

    # Salvar o arquivo PPTX
    nome_pptx = os.path.splitext(nome_arquivo_txt)[0] + ".pptx"
    prs.save(nome_pptx)
    # Converter PPTX para PDF usando LibreOffice 
    def pptx_para_pdf(input_pptx, output_pdf):
        sistema = platform.system()
        if sistema == "Windows":
            caminho_libreoffice = r"C:\Program Files\LibreOffice\program\soffice.exe"
            if not os.path.exists(caminho_libreoffice):
                raise FileNotFoundError(f"LibreOffice não encontrado no caminho especificado: {caminho_libreoffice}")
        elif sistema == "Darwin":  # macOS
            caminho_libreoffice = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
            if not os.path.exists(caminho_libreoffice):
                raise FileNotFoundError(f"LibreOffice não encontrado no caminho especificado: {caminho_libreoffice}")
        else:  # Linux e outros sistemas Unix-like
            caminho_libreoffice = "soffice"
            if subprocess.call(['which', 'soffice'], stdout=subprocess.PIPE, stderr=subprocess.PIPE) != 0:
                raise FileNotFoundError("LibreOffice não está instalado ou não está no PATH.")

        # Executar o comando para conversão
        comando = [
            caminho_libreoffice, "--headless", "--convert-to", "pdf",
            "--outdir", os.path.dirname(output_pdf), input_pptx
        ]
        subprocess.run(comando, check=True)

        print(f"Arquivo PDF gerado com sucesso: {output_pdf}")
    nome_pdf = os.path.splitext(nome_arquivo_txt)[0] + ".pdf"
    pptx_para_pdf(nome_pptx, nome_pdf)

# Listar todos os arquivos txt na pasta atual
for arquivo in os.listdir('.'):
    if arquivo.endswith('.txt'):
        criar_slides_a_partir_de_txt(arquivo)
